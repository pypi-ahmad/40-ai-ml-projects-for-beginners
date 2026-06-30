"""Report generation (Markdown, HTML, PDF, Excel, PowerPoint)."""

from __future__ import annotations

from pathlib import Path
from typing import Any

import markdown as md
import pandas as pd
from fpdf import FPDF
from jinja2 import Template
from loguru import logger
from pptx import Presentation

from ai_spreadsheet_analytics.schemas import ReportArtifacts
from ai_spreadsheet_analytics.utils import utc_timestamp


class ReportGenerator:
    """Generate downloadable business reports."""

    def __init__(self, report_dir: Path) -> None:
        self.report_dir = report_dir
        self.report_dir.mkdir(parents=True, exist_ok=True)

    def generate(
        self,
        title: str,
        insights_markdown: str,
        tables: dict[str, pd.DataFrame],
        formats: list[str] | None = None,
    ) -> ReportArtifacts:
        """Generate reports in selected formats."""
        requested = formats or ["markdown", "html", "pdf", "excel", "powerpoint"]
        slug = title.lower().replace(" ", "_")
        timestamp = utc_timestamp()
        base_name = f"{slug}_{timestamp}"

        artifacts = ReportArtifacts()
        markdown_body = self._build_markdown_content(title, insights_markdown, tables)

        if "markdown" in requested:
            try:
                path = self.report_dir / f"{base_name}.md"
                path.write_text(markdown_body, encoding="utf-8")
                artifacts.markdown = path
            except Exception as exc:  # noqa: BLE001
                logger.warning("Markdown report generation failed: {}", exc)

        if "html" in requested:
            try:
                path = self.report_dir / f"{base_name}.html"
                self._write_html(title, markdown_body, path)
                artifacts.html = path
            except Exception as exc:  # noqa: BLE001
                logger.warning("HTML report generation failed: {}", exc)

        if "pdf" in requested:
            try:
                path = self.report_dir / f"{base_name}.pdf"
                self._write_pdf(title, markdown_body, path)
                artifacts.pdf = path
            except Exception as exc:  # noqa: BLE001
                logger.warning("PDF report generation failed: {}", exc)

        if "excel" in requested:
            try:
                path = self.report_dir / f"{base_name}.xlsx"
                self._write_excel(tables, path)
                artifacts.excel = path
            except Exception as exc:  # noqa: BLE001
                logger.warning("Excel report generation failed: {}", exc)

        if "powerpoint" in requested:
            try:
                path = self.report_dir / f"{base_name}.pptx"
                self._write_powerpoint(title, insights_markdown, tables, path)
                artifacts.powerpoint = path
            except Exception as exc:  # noqa: BLE001
                logger.warning("PowerPoint report generation failed: {}", exc)

        return artifacts

    def _build_markdown_content(self, title: str, insights_markdown: str, tables: dict[str, pd.DataFrame]) -> str:
        lines = [f"# {title}", "", "## AI Insights", "", insights_markdown, "", "## Tables", ""]
        for name, table in tables.items():
            lines.append(f"### {name}")
            lines.append("")
            lines.append(table.head(20).to_markdown(index=False))
            lines.append("")
        return "\n".join(lines)

    def _write_html(self, title: str, markdown_body: str, output_path: Path) -> None:
        template = Template(
            """
            <html>
            <head><meta charset=\"utf-8\"><title>{{ title }}</title></head>
            <body style=\"font-family: Arial, sans-serif; max-width: 1100px; margin: 40px auto;\">
            {{ body | safe }}
            </body>
            </html>
            """
        )
        html_body = md.markdown(markdown_body, extensions=["tables"])
        output_path.write_text(template.render(title=title, body=html_body), encoding="utf-8")

    def _write_pdf(self, title: str, markdown_body: str, output_path: Path) -> None:
        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=12)
        pdf.set_font("Helvetica", size=12)
        pdf.cell(0, 8, txt=title.encode("latin-1", errors="ignore").decode("latin-1"), new_x="LMARGIN", new_y="NEXT")
        pdf.ln(2)

        lines: list[str] = ["Generated summary:"]
        for line in markdown_body.splitlines():
            stripped = line.strip()
            if not stripped:
                continue
            if stripped.startswith("|"):
                continue
            if stripped.startswith("#"):
                lines.append(stripped.lstrip("# ").strip())
                continue
            lines.append(stripped[:180])
            if len(lines) >= 180:
                break

        for line in lines:
            safe_line = line.encode("latin-1", errors="ignore").decode("latin-1")
            safe_line = "".join(ch for ch in safe_line if ch.isprintable())
            if not safe_line:
                continue
            pdf.cell(0, 6, txt=safe_line, new_x="LMARGIN", new_y="NEXT")
        pdf.output(str(output_path))

    def _write_excel(self, tables: dict[str, pd.DataFrame], output_path: Path) -> None:
        with pd.ExcelWriter(output_path, engine="xlsxwriter") as writer:
            for name, table in tables.items():
                table.to_excel(writer, sheet_name=name[:30], index=False)

    def _write_powerpoint(
        self,
        title: str,
        insights_markdown: str,
        tables: dict[str, pd.DataFrame],
        output_path: Path,
    ) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = insights_markdown[:2000]

        for name, table in tables.items():
            data_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            data_slide.shapes.title.text = name
            preview = table.head(8).to_string(index=False)
            data_slide.placeholders[1].text = preview[:4000]

        presentation.save(str(output_path))
