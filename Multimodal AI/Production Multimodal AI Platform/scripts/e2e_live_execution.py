"""Real end-to-end execution runner for multimodal AI platform."""

from __future__ import annotations

import json
import subprocess
import time
from pathlib import Path
from typing import Any

import httpx
from PIL import Image, ImageDraw


def create_sample_assets(base_dir: Path) -> dict[str, str]:
    """Create deterministic sample files for live e2e run."""

    base_dir.mkdir(parents=True, exist_ok=True)

    image_a = base_dir / "e2e_image_a.png"
    image_b = base_dir / "e2e_image_b.png"
    docx_path = base_dir / "e2e_notes.docx"
    pptx_path = base_dir / "e2e_slides.pptx"

    img = Image.new("RGB", (720, 420), color=(245, 247, 250))
    draw = ImageDraw.Draw(img)
    draw.rectangle((40, 40, 340, 220), outline=(30, 90, 140), width=4)
    draw.text((60, 70), "Quarterly Dashboard", fill=(30, 90, 140))
    draw.text((60, 110), "Revenue: +18%", fill=(20, 20, 20))
    draw.text((60, 145), "Errors: -12%", fill=(20, 20, 20))
    img.save(image_a)

    img2 = Image.new("RGB", (720, 420), color=(238, 244, 252))
    draw2 = ImageDraw.Draw(img2)
    draw2.rectangle((48, 50, 348, 230), outline=(20, 70, 120), width=4)
    draw2.text((68, 80), "Quarterly Dashboard", fill=(20, 70, 120))
    draw2.text((68, 120), "Revenue: +20%", fill=(20, 20, 20))
    draw2.text((68, 155), "Errors: -10%", fill=(20, 20, 20))
    img2.save(image_b)

    from docx import Document

    document = Document()
    document.add_heading("E2E Multimodal Report", level=1)
    document.add_paragraph("System processed chart and screenshot analytics for Q2 operations.")
    document.add_paragraph(
        "Primary insights: revenue increased, "
        "error rate decreased, latency stable."
    )
    document.save(docx_path)

    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Platform KPIs"
    slide.placeholders[1].text = "Availability 99.95%\nLatency P95: 220 ms\nTickets reduced by 14%"
    presentation.save(pptx_path)

    return {
        "image_a": str(image_a),
        "image_b": str(image_b),
        "docx": str(docx_path),
        "pptx": str(pptx_path),
        "pdf": str(base_dir / "example.pdf"),
    }


def start_process(cmd: list[str], log_path: Path) -> subprocess.Popen[bytes]:
    """Start background process with logs redirected to file."""

    log_path.parent.mkdir(parents=True, exist_ok=True)
    log_handle = log_path.open("wb")
    process = subprocess.Popen(cmd, stdout=log_handle, stderr=log_handle)
    return process


def wait_http(url: str, timeout_s: float = 120.0) -> str:
    """Wait for url reachable, return short response text."""

    start = time.time()
    last_error = ""
    while time.time() - start < timeout_s:
        try:
            response = httpx.get(url, timeout=5.0)
            response.raise_for_status()
            return response.text[:500]
        except Exception as exc:  # noqa: BLE001
            last_error = str(exc)
            time.sleep(0.7)
    raise RuntimeError(f"timeout waiting for {url}: {last_error}")


def run_cli(commands: list[list[str]]) -> list[dict[str, Any]]:
    """Run CLI command list and collect outputs."""

    records: list[dict[str, Any]] = []
    for command in commands:
        result = subprocess.run(command, capture_output=True, text=True, check=False)
        record = {
            "command": command,
            "exit_code": result.returncode,
            "stdout": result.stdout[-5000:],
            "stderr": result.stderr[-5000:],
        }
        if result.returncode != 0:
            raise RuntimeError(f"CLI failed: {command}\n{result.stdout}\n{result.stderr}")
        records.append(record)
    return records


def post_json(base_url: str, path: str, payload: dict[str, Any]) -> dict[str, Any]:
    """POST JSON and assert success."""

    response = httpx.post(f"{base_url}{path}", json=payload, timeout=120.0)
    if response.status_code != 200:
        raise RuntimeError(f"endpoint {path} failed: {response.status_code} {response.text}")
    body = response.json()
    if body.get("status") != "ok":
        raise RuntimeError(f"endpoint {path} returned non-ok: {body}")
    return body


def main() -> None:
    root = Path(".").resolve()
    artifacts_dir = root / "outputs" / "reports"
    artifacts_dir.mkdir(parents=True, exist_ok=True)
    api_port = 8011
    streamlit_port = 8512

    assets = create_sample_assets(root / "data" / "uploads")

    api_log = root / "outputs" / "reports" / "api_live.log"
    ui_log = root / "outputs" / "reports" / "streamlit_live.log"

    api_proc = start_process(
        [
            "uv",
            "run",
            "uvicorn",
            "app:app",
            "--host",
            "127.0.0.1",
            "--port",
            str(api_port),
        ],
        api_log,
    )
    ui_proc = start_process(
        [
            "uv",
            "run",
            "streamlit",
            "run",
            "src/multimodal_ai/ui/streamlit_app.py",
            "--server.address",
            "127.0.0.1",
            "--server.port",
            str(streamlit_port),
            "--server.headless",
            "true",
        ],
        ui_log,
    )

    try:
        api_health_raw = wait_http(f"http://127.0.0.1:{api_port}/health")
        streamlit_health_raw = wait_http(f"http://127.0.0.1:{streamlit_port}/_stcore/health")

        openapi = httpx.get(f"http://127.0.0.1:{api_port}/openapi.json", timeout=20.0).json()
        if "/caption" not in openapi.get("paths", {}):
            raise RuntimeError("live API openapi missing /caption route")

        base_url = f"http://127.0.0.1:{api_port}"
        trace = {"trace": {"source": "e2e_live"}}

        responses: dict[str, Any] = {}
        responses["caption"] = post_json(
            base_url,
            "/caption",
            {
                **trace,
                "input": {"image_path": assets["image_a"]},
                "options": {"style": "detailed"},
            },
        )
        responses["ocr_image"] = post_json(
            base_url,
            "/ocr",
            {
                **trace,
                "input": {"image_path": assets["image_a"]},
            },
        )
        responses["ocr_docx"] = post_json(
            base_url,
            "/ocr",
            {
                **trace,
                "input": {"document_path": assets["docx"]},
            },
        )
        responses["documents"] = post_json(
            base_url,
            "/documents",
            {
                **trace,
                "input": {"document_path": assets["docx"]},
            },
        )
        responses["embeddings"] = post_json(
            base_url,
            "/embeddings",
            {
                **trace,
                "input": {"text": "Find operational dashboard insights"},
            },
        )
        responses["compare"] = post_json(
            base_url,
            "/compare",
            {
                **trace,
                "input": {"image_paths": [assets["image_a"], assets["image_b"]]},
            },
        )
        responses["vqa"] = post_json(
            base_url,
            "/vqa",
            {
                **trace,
                "input": {
                    "document_path": assets["docx"],
                    "question": "What are key insights from this document?",
                },
            },
        )
        responses["search"] = post_json(
            base_url,
            "/search",
            {
                **trace,
                "input": {"query": "revenue increased"},
                "options": {"modality": "document", "top_k": 5},
            },
        )
        responses["retrieve"] = post_json(
            base_url,
            "/retrieve",
            {
                **trace,
                "input": {"query": "latency stable"},
                "options": {"modality": "document", "top_k": 5},
            },
        )
        responses["analyze"] = post_json(
            base_url,
            "/analyze",
            {
                **trace,
                "input": {
                    "image_path": assets["image_a"],
                    "question": "Explain this dashboard",
                },
            },
        )
        responses["analytics"] = post_json(base_url, "/analytics", trace)

        cli_results = run_cli(
            [
                ["uv", "run", "multimodal-ai", "doctor"],
                ["uv", "run", "multimodal-ai", "caption", assets["image_a"], "--style", "short"],
                ["uv", "run", "multimodal-ai", "ocr", assets["docx"]],
                ["uv", "run", "multimodal-ai", "analyze", "--image", assets["image_a"]],
            ]
        )

        summary = {
            "api_health_raw": api_health_raw,
            "streamlit_health_raw": streamlit_health_raw,
            "assets": assets,
            "responses": responses,
            "cli_results": cli_results,
            "api_log": str(api_log),
            "streamlit_log": str(ui_log),
        }

        json_path = artifacts_dir / "e2e_live_summary.json"
        md_path = artifacts_dir / "e2e_live_summary.md"

        json_path.write_text(json.dumps(summary, indent=2, default=str), encoding="utf-8")

        md_lines = [
            "# E2E Live Execution Summary",
            "",
            "## Services",
            f"- FastAPI health raw: `{api_health_raw[:120]}`",
            f"- Streamlit health raw: `{streamlit_health_raw[:120]}`",
            "",
            "## Endpoint Results",
        ]
        for key, value in responses.items():
            md_lines.append(
                f"- `{key}`: status={value.get('status')} "
                f"latency_ms={value.get('latency_ms')}"
            )
        md_lines.append("")
        md_lines.append("## CLI Results")
        for row in cli_results:
            md_lines.append(f"- `{' '.join(row['command'])}` -> exit_code={row['exit_code']}")
        md_path.write_text("\n".join(md_lines), encoding="utf-8")

        print(f"Wrote {json_path}")
        print(f"Wrote {md_path}")

    finally:
        for process in (api_proc, ui_proc):
            process.terminate()
            try:
                process.wait(timeout=10)
            except Exception:  # noqa: BLE001
                process.kill()


if __name__ == "__main__":
    main()
