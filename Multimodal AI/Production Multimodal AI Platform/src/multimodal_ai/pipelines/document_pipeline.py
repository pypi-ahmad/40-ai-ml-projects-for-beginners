"""Document parsing and OCR routing pipeline."""

from __future__ import annotations

import tempfile
from pathlib import Path

from multimodal_ai.adapters.ocr import estimate_scanned_document
from multimodal_ai.adapters.registry import AdapterRegistry
from multimodal_ai.domain import OCRResult


class DocumentPipeline:
    """Pipeline for document/text extraction with OCR fallback."""

    def __init__(
        self, registry: AdapterRegistry, primary_engine: str, min_text_chars: int = 40
    ) -> None:
        self._registry = registry
        self._primary_engine = primary_engine
        self._min_text_chars = min_text_chars

    def _extract_native_text(self, path: Path) -> str:
        ext = path.suffix.lower()
        if ext == ".pdf":
            return self._extract_pdf_text(path)
        if ext == ".docx":
            return self._extract_docx_text(path)
        if ext == ".pptx":
            return self._extract_pptx_text(path)
        return ""

    def _extract_pdf_text(self, path: Path) -> str:
        try:
            import pdfplumber

            chunks: list[str] = []
            with pdfplumber.open(path) as pdf:
                for page in pdf.pages:
                    chunks.append(page.extract_text() or "")
            return "\n".join(chunks).strip()
        except Exception:  # noqa: BLE001
            return ""

    def _extract_docx_text(self, path: Path) -> str:
        try:
            from docx import Document

            document = Document(str(path))
            return "\n".join(paragraph.text for paragraph in document.paragraphs).strip()
        except Exception:  # noqa: BLE001
            return ""

    def _extract_pptx_text(self, path: Path) -> str:
        try:
            from pptx import Presentation

            presentation = Presentation(str(path))
            lines: list[str] = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        lines.append(shape.text)
            return "\n".join(lines).strip()
        except Exception:  # noqa: BLE001
            return ""

    def _render_pdf_first_page(self, path: Path) -> Path | None:
        try:
            import fitz

            with fitz.open(path) as doc:
                page = doc[0]
                pix = page.get_pixmap()
                with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
                    pix.save(tmp.name)
                    return Path(tmp.name)
        except Exception:  # noqa: BLE001
            return None

    def run(self, file_path: str) -> OCRResult:
        """Run hybrid extraction for image or document."""

        path = Path(file_path)
        ext = path.suffix.lower()

        if ext in {".png", ".jpg", ".jpeg", ".webp", ".tiff"}:
            result = self._registry.create_ocr(self._primary_engine).extract(str(path))
            return OCRResult(
                engine=result.get("engine", self._primary_engine),
                text=result.get("text", ""),
                blocks=result.get("blocks", []),
                tables=result.get("tables", []),
                is_scanned=estimate_scanned_document(str(path)),
            )

        native_text = self._extract_native_text(path)
        if len(native_text) >= self._min_text_chars:
            return OCRResult(
                engine="native_parser",
                text=native_text,
                blocks=[],
                tables=[],
                is_scanned=False,
            )

        ocr_input = path
        if ext == ".pdf":
            rendered = self._render_pdf_first_page(path)
            if rendered is not None:
                ocr_input = rendered

        ocr_result = self._registry.create_ocr(self._primary_engine).extract(str(ocr_input))
        return OCRResult(
            engine=ocr_result.get("engine", self._primary_engine),
            text=ocr_result.get("text", ""),
            blocks=ocr_result.get("blocks", []),
            tables=ocr_result.get("tables", []),
            is_scanned=True,
        )
